"""
PtoF - PPTX to Figures - Core Engine

Module providing core processing logic.
"""

import os
import re
import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Color name to RGB value mapping
COLOR_NAMES = {
    'cyan': (0, 255, 255),
    'red': (255, 0, 0),
    'green': (0, 255, 0),
    'blue': (0, 0, 255),
    'magenta': (255, 0, 255),
    'yellow': (255, 255, 0),
    'orange': (255, 165, 0),
    'purple': (128, 0, 128),
    'pink': (255, 192, 203),
    'lime': (0, 255, 0),
    'black': (0, 0, 0),
    'white': (255, 255, 255),
}


def parse_color(color_str):
    """
    Convert a color string to an RGB tuple.

    Args:
        color_str: Color name ('cyan') or HEX format ('#00FFFF')

    Returns:
        tuple: (R, G, B)

    Raises:
        ValueError: If the color specification is invalid
    """
    color_str = color_str.strip().lower()

    # Color name
    if color_str in COLOR_NAMES:
        return COLOR_NAMES[color_str]

    # HEX format
    if color_str.startswith('#'):
        hex_str = color_str[1:]
        if len(hex_str) == 6:
            try:
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
                return (r, g, b)
            except ValueError:
                pass
        elif len(hex_str) == 3:
            try:
                r = int(hex_str[0] * 2, 16)
                g = int(hex_str[1] * 2, 16)
                b = int(hex_str[2] * 2, 16)
                return (r, g, b)
            except ValueError:
                pass

    available = ', '.join(COLOR_NAMES.keys())
    raise ValueError(f"Invalid color: '{color_str}'. Use color name ({available}) or HEX (#RRGGBB)")


def is_matching_color(color, target_rgb, tolerance=30):
    """
    Determine if a color is close to the target color.

    Args:
        color: RGBColor object or None
        target_rgb: Target color (R, G, B) tuple
        tolerance: Tolerance (allowed difference for each channel)

    Returns:
        bool: True if the color is close to the target
    """
    if color is None:
        return False

    try:
        r = color[0] if hasattr(color, '__getitem__') else color.red
        g = color[1] if hasattr(color, '__getitem__') else color.green
        b = color[2] if hasattr(color, '__getitem__') else color.blue
    except (TypeError, AttributeError):
        return False

    return (abs(r - target_rgb[0]) <= tolerance and
            abs(g - target_rgb[1]) <= tolerance and
            abs(b - target_rgb[2]) <= tolerance)


def get_shape_line_color(shape):
    """
    Get the line color of a shape.

    Args:
        shape: python-pptx Shape object

    Returns:
        RGBColor or None
    """
    try:
        line = shape.line
        if line.fill.type is not None:
            color = line.color.rgb
            return color
    except (AttributeError, TypeError):
        pass
    return None


def find_marker_rectangles(slide, target_rgb):
    """
    Detect rectangles with the specified color in a slide.

    Args:
        slide: python-pptx Slide object
        target_rgb: Target color (R, G, B) tuple

    Returns:
        list: List of rectangle info [{left, top, width, height, shape}] (in EMU)
    """
    rectangles = []

    for shape in slide.shapes:
        # Get line color
        line_color = get_shape_line_color(shape)

        if is_matching_color(line_color, target_rgb):
            rect = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'shape': shape,
            }
            rectangles.append(rect)

    return rectangles


def find_filename_textboxes(slide):
    """
    Detect text boxes containing "filename=" in a slide.

    Args:
        slide: python-pptx Slide object

    Returns:
        list: List of filename info [{filename, left, top, shape}]
    """
    filenames = []
    pattern = re.compile(r'filename\s*=\s*(\S+\.(?:pdf|png|svg))', re.IGNORECASE)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text
        match = pattern.search(text)

        if match:
            filenames.append({
                'filename': match.group(1),
                'left': shape.left,
                'top': shape.top,
                'shape': shape,
            })

    return filenames


def get_center(item):
    """
    Get the center coordinates of a shape.

    Args:
        item: dict with left, top, width, height

    Returns:
        tuple: (center_x, center_y)
    """
    cx = item['left'] + item.get('width', 0) / 2
    cy = item['top'] + item.get('height', 0) / 2
    return (cx, cy)


def calc_distance(item1, item2):
    """
    Calculate the distance between the centers of two shapes.

    Args:
        item1, item2: dict with left, top, width, height

    Returns:
        float: Euclidean distance between centers
    """
    c1 = get_center(item1)
    c2 = get_center(item2)
    return ((c1[0] - c2[0]) ** 2 + (c1[1] - c2[1]) ** 2) ** 0.5


def match_rectangles_to_filenames(rectangles, filenames):
    """
    Match rectangles to filename text boxes based on distance.

    Args:
        rectangles: List of rectangle info
        filenames: List of filename info

    Returns:
        list: List of matched pairs [(rect, filename_info), ...]
    """
    if not rectangles or not filenames:
        return []

    # Calculate distances for all pairs
    pairs = []
    for rect in rectangles:
        for fn in filenames:
            dist = calc_distance(rect, fn)
            pairs.append((dist, rect, fn))

    # Sort by distance
    pairs.sort(key=lambda x: x[0])

    # Greedy matching (each rectangle and filename used only once)
    matched = []
    used_rects = set()
    used_filenames = set()

    for dist, rect, fn in pairs:
        rect_id = id(rect)
        fn_id = id(fn)

        if rect_id not in used_rects and fn_id not in used_filenames:
            matched.append((rect, fn))
            used_rects.add(rect_id)
            used_filenames.add(fn_id)

    return matched


def convert_pptx_to_pdf(pptx_path, pdf_path, embed_fonts=False):
    """
    Convert PPTX to PDF using PowerPoint COM.

    Args:
        pptx_path: Input PPTX file path
        pdf_path: Output PDF file path
        embed_fonts: True to force font embedding (PDF/A format)
    """
    import win32com.client

    powerpoint = None
    presentation = None

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")

        # Convert to absolute path
        pptx_abs = str(Path(pptx_path).resolve())
        pdf_abs = str(Path(pdf_path).resolve())

        # Open presentation (ReadOnly=True, WithWindow=False)
        presentation = powerpoint.Presentations.Open(pptx_abs, True, False, False)

        # Save as PDF
        # ppFixedFormatTypePDF = 2
        # ppFixedFormatIntentPrint = 2 (high quality)
        if embed_fonts:
            # PDF/A format (ISO 19005-1) with forced font embedding
            presentation.ExportAsFixedFormat(
                pdf_abs,
                2,  # ppFixedFormatTypePDF
                Intent=2,  # ppFixedFormatIntentPrint
                FrameSlides=False,
                HandoutOrder=1,
                OutputType=1,
                PrintHiddenSlides=False,
                IncludeDocProperties=True,
                KeepIRMSettings=True,
                DocStructureTags=True,
                BitmapMissingFonts=True,
                UseISO19005_1=True  # PDF/A format
            )
        else:
            presentation.SaveAs(pdf_abs, 32)  # 32 = ppSaveAsPDF

    finally:
        if presentation:
            presentation.Close()
        if powerpoint:
            powerpoint.Quit()


def clip_region(input_pdf_path, output_path, page_num, rect, slide_width, slide_height, dpi=300):
    """
    Clip a specific region from a PDF page.

    Args:
        input_pdf_path: Input PDF file path
        output_path: Output file path (.pdf, .png, .svg)
        page_num: Page number (0-indexed)
        rect: Clipping region {left, top, width, height} (in EMU)
        slide_width: Slide width (in EMU)
        slide_height: Slide height (in EMU)
        dpi: Resolution for PNG output (default: 300)
    """
    import fitz  # PyMuPDF

    doc = fitz.open(input_pdf_path)
    page = doc[page_num]

    # Get actual PDF page size
    pdf_width = page.rect.width
    pdf_height = page.rect.height

    # Convert EMU to PDF coordinates
    scale_x = pdf_width / slide_width
    scale_y = pdf_height / slide_height

    x0 = rect['left'] * scale_x
    y0 = rect['top'] * scale_y
    x1 = x0 + rect['width'] * scale_x
    y1 = y0 + rect['height'] * scale_y

    # Set clipping region
    clip_rect = fitz.Rect(x0, y0, x1, y1)

    # Determine output format
    output_path = str(output_path)
    ext = output_path.lower().rsplit('.', 1)[-1]

    if ext == 'png':
        # Output as PNG
        zoom = dpi / 72  # 72 DPI is the baseline
        mat = fitz.Matrix(zoom, zoom)
        pixmap = page.get_pixmap(matrix=mat, clip=clip_rect)
        pixmap.save(output_path)

    elif ext == 'svg':
        # Output as SVG (create a clipped PDF page first, then convert)
        temp_doc = fitz.open()
        temp_page = temp_doc.new_page(width=clip_rect.width, height=clip_rect.height)
        temp_page.show_pdf_page(temp_page.rect, doc, page_num, clip=clip_rect)
        svg_content = temp_page.get_svg_image()
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(svg_content)
        temp_doc.close()

    else:
        # Output as PDF
        new_doc = fitz.open()
        new_page = new_doc.new_page(width=clip_rect.width, height=clip_rect.height)

        # Copy clipping region from original page
        new_page.show_pdf_page(
            new_page.rect,
            doc,
            page_num,
            clip=clip_rect
        )

        new_doc.save(output_path)
        new_doc.close()

    doc.close()


def remove_shape(shape):
    """
    Remove a shape from the slide.

    Args:
        shape: Shape to remove
    """
    sp = shape._element
    sp.getparent().remove(sp)


def scan_pptx(pptx_path, marker_color=(0, 255, 255)):
    """
    Scan a PPTX file to get clipping information.

    Args:
        pptx_path: Input PPTX file path
        marker_color: Marker rectangle color (R, G, B) tuple

    Returns:
        dict: {
            'slide_width': Slide width,
            'slide_height': Slide height,
            'clip_info': List of clipping info,
            'shapes_to_remove': List of shapes to remove,
            'presentation': Presentation object
        }
    """
    pptx_path = Path(pptx_path)
    base_name = pptx_path.stem

    # Load PPTX
    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Collect clipping info from each slide
    clip_info = []
    shapes_to_remove = []

    for slide_idx, slide in enumerate(prs.slides):
        rectangles = find_marker_rectangles(slide, marker_color)
        filenames = find_filename_textboxes(slide)

        if not rectangles:
            continue

        # Match rectangles to filenames based on distance
        matched_pairs = match_rectangles_to_filenames(rectangles, filenames)

        # Record IDs of matched rectangles
        matched_rect_ids = {id(rect) for rect, _ in matched_pairs}

        # Process matched rectangles
        for rect, filename_info in matched_pairs:
            clip_info.append({
                'slide_idx': slide_idx,
                'rect': {
                    'left': rect['left'],
                    'top': rect['top'],
                    'width': rect['width'],
                    'height': rect['height'],
                },
                'filename': filename_info['filename'],
            })

            shapes_to_remove.append(rect['shape'])
            shapes_to_remove.append(filename_info['shape'])

        # Assign default names to unmatched rectangles
        unmatched_idx = 0
        for rect in rectangles:
            if id(rect) not in matched_rect_ids:
                unmatched_idx += 1
                default_filename = f"{base_name}_s{slide_idx + 1}_{len(matched_pairs) + unmatched_idx}.pdf"
                clip_info.append({
                    'slide_idx': slide_idx,
                    'rect': {
                        'left': rect['left'],
                        'top': rect['top'],
                        'width': rect['width'],
                        'height': rect['height'],
                    },
                    'filename': default_filename,
                })
                shapes_to_remove.append(rect['shape'])

    return {
        'slide_width': slide_width,
        'slide_height': slide_height,
        'clip_info': clip_info,
        'shapes_to_remove': shapes_to_remove,
        'presentation': prs,
    }


def process_pptx(pptx_path, output_dir, embed_fonts=False, marker_color=(0, 255, 255),
                 dpi=300, margin=0, dry_run=False, quiet=False, no_overwrite=False,
                 progress_callback=None):
    """
    Process a PPTX file and output clipped PDFs.

    Args:
        pptx_path: Input PPTX file path
        output_dir: Output directory path
        embed_fonts: True to force font embedding
        marker_color: Marker rectangle color (R, G, B) tuple
        dpi: Resolution for PNG output
        margin: Margin in points (positive: expand, negative: shrink)
        dry_run: If True, only show detection results
        quiet: If True, suppress output
        no_overwrite: If True, confirm before overwriting existing files
        progress_callback: Progress callback function (message, current, total) -> bool
                          Returns False to cancel

    Returns:
        list: List of output files
    """
    def log(msg):
        if not quiet:
            print(msg)

    def progress(msg, current=0, total=0):
        if progress_callback:
            return progress_callback(msg, current, total)
        return True

    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Scan
    scan_result = scan_pptx(pptx_path, marker_color)
    clip_info = scan_result['clip_info']
    shapes_to_remove = scan_result['shapes_to_remove']
    prs = scan_result['presentation']
    slide_width = scan_result['slide_width']
    slide_height = scan_result['slide_height']

    if not clip_info:
        log("No clipping regions found in the presentation")
        return []

    # Convert margin to EMU (1 point = 12700 EMU)
    margin_emu = int(margin * 12700)

    # Apply margin
    for info in clip_info:
        info['rect']['left'] -= margin_emu
        info['rect']['top'] -= margin_emu
        info['rect']['width'] += margin_emu * 2
        info['rect']['height'] += margin_emu * 2

    # In dry-run mode, show detection results and exit
    if dry_run:
        log(f"\n[Dry-run] {pptx_path.name}:")
        for info in clip_info:
            log(f"  Slide {info['slide_idx'] + 1} -> {info['filename']}")
        return [output_dir / info['filename'] for info in clip_info]

    # Check for existing files
    if no_overwrite:
        existing_files = []
        for info in clip_info:
            output_path = output_dir / info['filename']
            if output_path.exists():
                existing_files.append(output_path)

        if existing_files:
            print("The following files already exist:")
            for f in existing_files:
                print(f"  - {f}")
            response = input("Overwrite? [y/N]: ").strip().lower()
            if response != 'y':
                print("Aborted.")
                return []

    # Remove marker rectangles and filename text boxes
    for shape in shapes_to_remove:
        remove_shape(shape)

    # Create temporary PPTX and PDF files
    temp_dir = tempfile.mkdtemp()
    temp_pptx_path = os.path.join(temp_dir, 'temp.pptx')
    temp_pdf_path = os.path.join(temp_dir, 'temp.pdf')

    try:
        # Save modified PPTX
        prs.save(temp_pptx_path)

        log(f"Converting {pptx_path.name} to PDF...")
        if not progress(f"Converting {pptx_path.name} to PDF...", 0, len(clip_info)):
            return []

        convert_pptx_to_pdf(temp_pptx_path, temp_pdf_path, embed_fonts)

        output_files = []

        # Process each clipping region
        for i, info in enumerate(clip_info):
            output_path = output_dir / info['filename']
            log(f"Clipping slide {info['slide_idx'] + 1} -> {info['filename']}")

            if not progress(f"Clipping {info['filename']}", i + 1, len(clip_info)):
                return output_files

            clip_region(
                temp_pdf_path,
                str(output_path),
                info['slide_idx'],
                info['rect'],
                slide_width,
                slide_height,
                dpi
            )

            output_files.append(output_path)

        return output_files

    finally:
        # Delete temporary files
        shutil.rmtree(temp_dir, ignore_errors=True)
